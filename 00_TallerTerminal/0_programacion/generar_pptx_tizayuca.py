# Tizayuca — Generador de PPTX — Sistema Visual Industrial Organic
# Basado en analisis-urbano-v1.html
# Ejecutar en Google Colab

# ─── DEPENDENCIAS ────────────────────────────────────────────
!pip install python-pptx requests -q

import requests, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

# ─── DATOS ───────────────────────────────────────────────────
GITHUB_RAW = (
    "https://raw.githubusercontent.com/"
    "guillermojhoel2718-code/Taller-Terminal-ESIA/main/"
    "00_TallerTerminal/01_Investigacion/Datos/Municipios/02_Tizayuca.json"
)
r = requests.get(GITHUB_RAW)
D = r.json()
print("✅ JSON cargado:", D["municipio"], D["estado"])

# ─── PALETA (del CSS :root del HTML) ─────────────────────────
BG_DARK  = RGBColor(0x1A, 0x1A, 0x1A)  # --bg-dark
BG_LIGHT = RGBColor(0xF5, 0xF0, 0xE8)  # --bg-light
OCRE     = RGBColor(0xC8, 0xA9, 0x6E)  # --ocre  (labels, hero title)
MUSGO    = RGBColor(0x4A, 0x7C, 0x59)  # --tiza  (Tizayuca)
TERRA    = RGBColor(0xD4, 0x44, 0x1C)  # --terra (alertas, badge versión)
ACERO    = RGBColor(0x6B, 0x8F, 0xA3)  # --acero
GRIS     = RGBColor(0x66, 0x66, 0x66)  # --gris
BLANCO   = RGBColor(0xFF, 0xFF, 0xFF)
AMARILLO = RGBColor(0xFF, 0xCC, 0x00)
GRIS_MID = RGBColor(0xAA, 0xAA, 0xAA)  # textos secundarios
BORDE_DK = RGBColor(0x2A, 0x2A, 0x2A)  # bordes en dark
BORDE_LT = RGBColor(0xE0, 0xD8, 0xCC)  # bordes en light

# ─── CONFIGURACIÓN ───────────────────────────────────────────
W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def fondo(slide, color=BG_DARK):
    """Fondo plano de slide."""
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def rect(slide, left, top, width, height, fill, line=None):
    """Forma rectangularcon relleno sólido, sin borde por defecto."""
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, left, top, width, height,
        size=14, bold=False, italic=False,
        color=BG_LIGHT, align=PP_ALIGN.LEFT,
        font="Calibri", wrap=True):
    """Textbox genérico."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text  = str(text)
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return tb

# — Etiqueta de sección (SEC-LABEL del HTML) —————————————
def sec_label(slide, texto, left=Inches(0.55), top=Inches(0.28)):
    """Texto en ocre, mayúsculas, letter-spacing visual, 9pt Calibri."""
    txt(slide, texto.upper(), left, top, Inches(12), Inches(0.35),
        size=8, bold=True, color=OCRE, font="Calibri")

# — Título de sección (h2 del HTML) ——————————————————————
def sec_h2(slide, texto, left=Inches(0.55), top=Inches(0.6)):
    txt(slide, texto, left, top, Inches(12), Inches(0.75),
        size=26, bold=True, color=BLANCO, font="Georgia")

# — Descripción corta de sección —————————————————————————
def sec_desc(slide, texto, left=Inches(0.55), top=Inches(1.3)):
    txt(slide, texto, left, top, Inches(12), Inches(0.55),
        size=11, color=GRIS_MID, font="Calibri")

# — Barra de acento musgo (izquierda) ————————————————————
def barra_musgo(slide):
    rect(slide, Inches(0.3), Inches(0.22), Inches(0.06), Inches(7.1), MUSGO)

# — Línea divisora horizontal ————————————————————————————
def h_line(slide, top, color=BORDE_DK):
    rect(slide, Inches(0.5), top, Inches(12.3), Pt(1), color)

# — Stat card (stat-card del HTML) ———————————————————————
def stat_card(slide, left, top, ancho, alto,
              mun_label, valor, label, sub="",
              accent=MUSGO, dark=True):
    """Card con borde superior de color, estilo HTML."""
    bg = RGBColor(0x22,0x22,0x22) if dark else BLANCO
    # fondo del card
    rect(slide, left, top, ancho, alto, bg, BORDE_DK if dark else BORDE_LT)
    # borde superior de acento (border-top: 3px solid)
    rect(slide, left, top, ancho, Pt(3), accent)
    # mun label — estilo .stat-mun
    txt(slide, mun_label.upper(), left+Inches(0.15), top+Inches(0.12),
        ancho-Inches(0.3), Inches(0.3),
        size=7, bold=True, color=RGBColor(0x99,0x99,0x99), font="Calibri")
    # número grande — estilo .stat-num
    txt(slide, str(valor), left+Inches(0.12), top+Inches(0.38),
        ancho-Inches(0.24), Inches(0.85),
        size=28, bold=True, color=accent, font="Georgia")
    # label — estilo .stat-label
    txt(slide, label, left+Inches(0.15), top+Inches(1.18),
        ancho-Inches(0.3), Inches(0.3),
        size=9, color=RGBColor(0xBB,0xBB,0xBB) if dark else GRIS, font="Calibri")
    # sub — estilo .stat-sub
    if sub:
        h_line(slide, top+Inches(1.5),
               BORDE_DK if dark else BORDE_LT)
        txt(slide, sub, left+Inches(0.15), top+Inches(1.55),
            ancho-Inches(0.3), Inches(0.5),
            size=8, color=GRIS_MID if dark else GRIS, font="Calibri")

# — Barra horizontal (bar chart del HTML) ————————————————
def bar_row(slide, label, valor_pct, fill_color,
            left, top, width=Inches(5.5), track_h=Inches(0.12)):
    txt(slide, label, left, top, width*0.6, Inches(0.28),
        size=12, color=fill_color, font="Calibri")
    txt(slide, f"{valor_pct}%", left+width*0.65, top, width*0.35, Inches(0.28),
        size=12, bold=True, color=BLANCO, align=PP_ALIGN.RIGHT, font="Calibri")
    # track
    rect(slide, left, top+Inches(0.3), width, track_h,
         RGBColor(0x30,0x30,0x30))
    # fill
    fill_w = int(width * (min(valor_pct,100)/100))
    if fill_w > 0:
        rect(slide, left, top+Inches(0.3), fill_w, track_h, fill_color)

# — Def-card (border-left del HTML) ——————————————————————
def def_card(slide, left, top, width, height,
             mun, title, body, accent=MUSGO, dark=False):
    bg = RGBColor(0x22,0x22,0x22) if dark else BLANCO
    b_color = BORDE_DK if dark else BORDE_LT
    rect(slide, left, top, width, height, bg, b_color)
    # borde izquierdo (border-left: 4px solid)
    rect(slide, left, top, Pt(4), height, accent)
    # label municipio
    txt(slide, mun.upper(), left+Inches(0.2), top+Inches(0.12),
        width-Inches(0.3), Inches(0.25),
        size=7, bold=True, color=accent, font="Calibri")
    # title
    txt(slide, title, left+Inches(0.2), top+Inches(0.35),
        width-Inches(0.3), Inches(0.35),
        size=11, bold=True,
        color=BG_DARK if not dark else BLANCO, font="Georgia")
    # body
    txt(slide, body, left+Inches(0.2), top+Inches(0.68),
        width-Inches(0.3), height-Inches(0.78),
        size=10, color=GRIS if not dark else GRIS_MID,
        font="Calibri", wrap=True)

# — Badge de versión (nav-version del HTML) ———————————————
def badge_version(slide, ver="v1.0", left=Inches(12.0), top=Inches(0.12)):
    rect(slide, left, top, Inches(0.95), Inches(0.28), TERRA)
    txt(slide, ver, left, top+Inches(0.02), Inches(0.95), Inches(0.28),
        size=8, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, font="Courier New")

# — Cita fuente pequeña ——————————————————————————————————
def fuente_pie(slide, texto, top=Inches(7.0)):
    txt(slide, f"Fuente: {texto}", Inches(0.5), top, Inches(12.3), Inches(0.35),
        size=7, italic=True, color=RGBColor(0x55,0x55,0x55), font="Calibri")

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA (hero del HTML)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
fondo(slide)

# Línea ocre superior (nav border-bottom adaptado)
rect(slide, 0, 0, W, Inches(0.06), OCRE)
# Línea ocre inferior
rect(slide, 0, H-Inches(0.06), W, Inches(0.06), OCRE)

# hero-tag
txt(slide, "TALLER TERMINAL I · INGENIERÍA ARQUITECTÓNICA",
    Inches(1), Inches(1.4), Inches(11), Inches(0.35),
    size=8, bold=True, color=OCRE, align=PP_ALIGN.CENTER, font="Calibri")

# h1 con span (Georgia bold)
txt(slide, "TIZAYUCA, HIDALGO",
    Inches(0.8), Inches(1.85), Inches(11.7), Inches(1.3),
    size=48, bold=True, color=MUSGO, align=PP_ALIGN.CENTER, font="Georgia")

txt(slide, "Centro de Teletrabajo y Aprendizaje Digital",
    Inches(1), Inches(3.1), Inches(11), Inches(0.6),
    size=20, color=BLANCO, align=PP_ALIGN.CENTER, font="Georgia")

txt(slide, "Protocolo de Investigación — Unidad I",
    Inches(1), Inches(3.75), Inches(11), Inches(0.45),
    size=13, italic=True, color=GRIS_MID, align=PP_ALIGN.CENTER, font="Calibri")

# Chip de municipio (chip-tiza del HTML)
chip_left = Inches(5.4)
rect(slide, chip_left, Inches(4.4), Inches(2.5), Inches(0.38), MUSGO)
txt(slide, "Tizayuca · Hidalgo · 45 km",
    chip_left+Inches(0.12), Inches(4.44), Inches(2.3), Inches(0.35),
    size=9, bold=True, color=BLANCO, font="Calibri")

# hero-meta (separador + datos)
h_line(slide, Inches(5.05), BORDE_DK)
metas = [
    ("ALUMNO", "Hernández Gómez Guillermo Jhoel"),
    ("CORREDOR", "MEX-85D Norte CDMX"),
    ("FUENTES", "INEGI 2020 · CONEVAL 2020 · SIGEH"),
    ("VERSIÓN", "v1.0 · Marzo 2026"),
]
for i,(lbl,val) in enumerate(metas):
    lx = Inches(0.8 + i*3.1)
    txt(slide, lbl, lx, Inches(5.2), Inches(2.9), Inches(0.25),
        size=7, color=RGBColor(0x55,0x55,0x55), font="Calibri")
    txt(slide, val, lx, Inches(5.45), Inches(2.9), Inches(0.4),
        size=10, color=OCRE, font="Calibri")

badge_version(slide)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — LOCALIZACIÓN (dark)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
fondo(slide)
barra_musgo(slide)

sec_label(slide, "01 · Localización Geográfica")
txt(slide, "Localización Geográfica", Inches(0.55), Inches(0.6),
    Inches(12), Inches(0.75), size=26, bold=True, color=BLANCO, font="Georgia")

reg = D.get("region_geografica", {})
coord = D.get("coordenadas", {})
col = D.get("colindancias", {})

desc = (f"México → Eje Neovolcánico → Hidalgo → Región 6 Tizayuca → "
        f"ZMVM (único municipio hidalguense) · {D['distancia_azcapo_km']} km por MEX-85D")
sec_desc(slide, desc)

# 4 stat cards: coordenadas, altitud, superficie, ZMVM
card_data = [
    ("Centroide", "19.8417° N\n98.9892° O", "Coordenadas municipales"),
    ("2,300 msnm", "", "Altitud oficial INEGI"),
    ("76.81 km²", "0.37% del estado", "Superficie municipal"),
    ("ZMVM", "Único municipio Hgo.", "Zona Metropolitana Valle de México"),
]
cw = Inches(2.8)
for i,(val,sub,lbl) in enumerate(card_data):
    stat_card(slide,
              Inches(0.5 + i*3.2), Inches(2.0),
              cw, Inches(2.2),
              "", val, lbl, sub, MUSGO)

# Colindancias
h_line(slide, Inches(4.45))
txt(slide, "COLINDANCIAS", Inches(0.55), Inches(4.6),
    Inches(12), Inches(0.25), size=7, bold=True, color=OCRE, font="Calibri")

col_items = [
    f"N: {col.get('norte','')}",
    f"S: {col.get('sur','')}",
    f"E: {col.get('este','')}  ·  O: {col.get('oeste','')}",
]
for i,c in enumerate(col_items):
    txt(slide, c, Inches(0.55), Inches(4.9+i*0.45), Inches(12), Inches(0.4),
        size=11, color=GRIS_MID, font="Calibri")

fuente_pie(slide, "INEGI, Compendio de información geográfica municipal 2010 · CONAPO/INEGI (ZMVM)")
badge_version(slide)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — MEDIO FÍSICO (dark)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
fondo(slide)
barra_musgo(slide)

sec_label(slide, "02 · Medio Físico")
txt(slide, "Medio Físico", Inches(0.55), Inches(0.6),
    Inches(12), Inches(0.75), size=26, bold=True, color=BLANCO, font="Georgia")
sec_desc(slide, "Topografía · Geología · Edafología · Clima · Vegetación y Fauna · Hidrografía")

geo = D.get("geologia", {})
eda = D.get("edafologia", {})
cli = D.get("clima", {})
veg = D.get("vegetacion_fauna", {})

# Columna izquierda — 4 def-cards apilados
card_items = [
    ("TOPOGRAFÍA",    D.get("topografia",""),
     f"Parte del Eje Neovolcánico · {geo.get('contexto','')}"),
    ("GEOLOGÍA",      geo.get("roca_dominante",""),
     f"Período: {geo.get('periodo_principal','')}"),
    ("EDAFOLOGÍA",    eda.get("suelo_dominante",""),
     f"Aptitud constructiva: moderada-baja · cimentaciones profundas requeridas"),
    ("CLIMA",         f"{cli.get('clasificacion_koppen','')} · {cli.get('temp_media_anual_c','')}°C media",
     f"Precip. {cli.get('precipitacion_media_mm','')} · Lluvia jun-oct"),
]

for i,(lbl,title,body) in enumerate(card_items):
    def_card(slide,
             Inches(0.5), Inches(1.75 + i*1.35),
             Inches(6.1), Inches(1.28),
             lbl, title, body, MUSGO, dark=True)

# Columna derecha — 2 cards grandes
def_card(slide,
         Inches(6.9), Inches(1.75), Inches(6.0), Inches(2.6),
         "VEGETACIÓN Y FAUNA",
         veg.get("vegetacion_original",""),
         f"Estado actual: {veg.get('estado_actual','')}\n\n"
         f"Fauna: {veg.get('fauna_grupos','')}\n"
         f"Especie en riesgo: {veg.get('especie_en_riesgo','')}",
         MUSGO, dark=True)

rie = D.get("riesgos_naturales", {})
def_card(slide,
         Inches(6.9), Inches(4.5), Inches(6.0), Inches(2.1),
         "RIESGOS NATURALES",
         "Subsidencia + Sismicidad moderada baja",
         rie.get("subsidencia",""),
         TERRA, dark=True)

fuente_pie(slide, "INEGI Compendio Municipal 2010 · Atlas de Riesgos Tizayuca 2013 · CENAPRED")
badge_version(slide)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — DATOS SOCIODEMOGRÁFICOS (dark, stat-grid)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
fondo(slide)
barra_musgo(slide)

sec_label(slide, "03 · Demografía")
txt(slide, "Datos Sociodemográficos", Inches(0.55), Inches(0.6),
    Inches(12), Inches(0.75), size=26, bold=True, color=BLANCO, font="Georgia")
sec_desc(slide, f"Población total 2020: {D['poblacion']:,} hab · "
                f"Crecimiento intercensal +72.7% · Densidad {D['densidad_hab_km2']:,} hab/km²")

# Fila 1 — 4 stat cards grandes
cards_row1 = [
    (f"{D['poblacion']:,}", "habitantes 2020", "Crecimiento +72.7%"),
    (f"{D['densidad_hab_km2']:,}", "hab/km²", "Mayor densidad del corredor"),
    (f"{D['viviendas_habitadas']:,}", "viviendas habitadas", f"{D['ocupantes_vivienda']} ocup./vivienda"),
    ("+72.7%", "crecimiento intercensal", "2010 → 2020 (de 97,461 hab)"),
]
cw = Inches(2.9)
for i,(val,lbl,sub) in enumerate(cards_row1):
    stat_card(slide, Inches(0.4+i*3.15), Inches(1.9), cw, Inches(2.0),
              "TIZAYUCA", val, lbl, sub, MUSGO)

# Fila 2 — barras de pobreza/rezago
h_line(slide, Inches(4.15))
txt(slide, "INDICADORES SOCIALES", Inches(0.55), Inches(4.25),
    Inches(12), Inches(0.25), size=7, bold=True, color=OCRE, font="Calibri")

bars = [
    ("Pobreza total",         float(D["pobreza_pct"]),       TERRA),
    ("Rezago educativo",      float(D["rezago_educativo_pct"]), ACERO),
    ("Internet en viviendas", float(D["internet_pct"]),       MUSGO),
    ("Analfabetismo",         float(D["analfabetismo_pct"]),  GRIS_MID),
]
for i,(lbl,val,col_b) in enumerate(bars):
    bar_row(slide, f"{lbl}  {val}%", val, col_b,
            Inches(0.5 + (i//2)*6.4),
            Inches(4.6 + (i%2)*0.82))

fuente_pie(slide, "INEGI Censo 2020 · CONEVAL Medición de Pobreza Municipal 2020")
badge_version(slide)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — SERVICIOS BÁSICOS (light-style sobre dark)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
fondo(slide)
barra_musgo(slide)

sec_label(slide, "04 · Infraestructura")
txt(slide, "Servicios Básicos", Inches(0.55), Inches(0.6),
    Inches(12), Inches(0.75), size=26, bold=True, color=BLANCO, font="Georgia")
sec_desc(slide, "Cobertura sobre total de viviendas habitadas · Fuente: INEGI Censo 2020")

servicios = [
    ("💧 Agua potable",    D["agua_pct"],        ACERO),
    ("🚿 Drenaje",         D["drenaje_pct"],      MUSGO),
    ("⚡ Electricidad",    D["electricidad_pct"], OCRE),
    ("🌐 Internet",        D["internet_pct"],     TERRA),
    ("💻 Computadora",     D["computadora_pct"],  GRIS_MID),
]

cols = 3
for i,(lbl,val,acent) in enumerate(servicios):
    col_i = i % cols
    row_i = i // cols
    lx = Inches(0.5 + col_i*4.2)
    ty = Inches(1.85 + row_i*2.8)
    card_w, card_h = Inches(3.8), Inches(2.5)
    # card fondo
    rect(slide, lx, ty, card_w, card_h, RGBColor(0x22,0x22,0x22), BORDE_DK)
    rect(slide, lx, ty, card_w, Pt(3), acent)
    # valor grande
    txt(slide, f"{val}%", lx, ty+Inches(0.2), card_w, Inches(1.3),
        size=44, bold=True, color=acent, align=PP_ALIGN.CENTER, font="Georgia")
    txt(slide, lbl, lx, ty+Inches(1.5), card_w, Inches(0.4),
        size=12, color=GRIS_MID, align=PP_ALIGN.CENTER, font="Calibri")

fuente_pie(slide, "INEGI. (2020). Censo de Población y Vivienda 2020.")
badge_version(slide)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — HIDROGRAFÍA Y RIESGOS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
fondo(slide)
barra_musgo(slide)

sec_label(slide, "05 · Hidrografía y Riesgos")
txt(slide, "Hidrografía y Riesgos Naturales", Inches(0.55), Inches(0.6),
    Inches(12), Inches(0.75), size=26, bold=True, color=BLANCO, font="Georgia")

# Bloque de hidrografía
def_card(slide, Inches(0.5), Inches(1.65), Inches(12.3), Inches(1.5),
         "HIDROGRAFÍA",
         "Región Pánuco · Cuenca R. Moctezuma · Subcuenca R. Tezontepec",
         f"{D.get('hidrografia','')}\nPresa El Manantial como cuerpo de agua principal",
         ACERO, dark=True)

# Bloque riesgos
rie = D.get("riesgos_naturales", {})
def_card(slide, Inches(0.5), Inches(3.35), Inches(6.0), Inches(2.8),
         "SUBSIDENCIA / INUNDACIÓN",
         "Susceptibilidad alta en planicie aluvial",
         rie.get("subsidencia",""),
         TERRA, dark=True)

def_card(slide, Inches(6.8), Inches(3.35), Inches(6.0), Inches(2.8),
         "SISMICIDAD",
         "Moderada baja — Eje Neovolcánico",
         rie.get("sismicidad","") + "\n\nFuente: " + rie.get("fuente_sismicidad",""),
         OCRE, dark=True)

fuente_pie(slide, "Atlas de Riesgos Naturales de Tizayuca, Hgo. 2013 · CENAPRED")
badge_version(slide)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — DÉFICIT Y JUSTIFICACIÓN + PENDIENTES
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
fondo(slide)
barra_musgo(slide)

sec_label(slide, "06 · Déficit Identificado")
txt(slide, "Déficit y Justificación del Proyecto", Inches(0.55), Inches(0.6),
    Inches(12), Inches(0.75), size=26, bold=True, color=BLANCO, font="Georgia")

# Box ocre (estilo data-box del HTML)
rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.1),
     RGBColor(0x20,0x18,0x08), OCRE)
txt(slide, "💡  Dato clave para el proyecto:",
    Inches(0.7), Inches(1.7), Inches(12), Inches(0.35),
    size=10, bold=True, color=OCRE, font="Calibri")
txt(slide, D.get("deficit_principal","") + "\n\n"
    "Tizayuca: 23% de trabajadores tarda más de 1h en trasladarse. "
    "Un Centro de Teletrabajo ahorraría 2–4h diarias por persona.",
    Inches(0.7), Inches(2.05), Inches(12), Inches(1.4),
    size=12, color=BLANCO, font="Calibri", wrap=True)

# Box terra — pendientes
pendientes = [k for k,v in D.items() if v == "[PENDIENTE]"]
if pendientes:
    rect(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.28),
         RGBColor(0x20,0x08,0x05), TERRA)
    txt(slide, "⚠️  DATOS PENDIENTES DE RECOLECCIÓN:",
        Inches(0.65), Inches(3.93), Inches(12), Inches(0.25),
        size=8, bold=True, color=TERRA, font="Calibri")
    txt(slide, "  ·  ".join(pendientes),
        Inches(0.65), Inches(4.25), Inches(12.1), Inches(1.5),
        size=9, color=RGBColor(0x88,0x88,0x88), font="Calibri", wrap=True)

fuente_pie(slide, "INEGI 2020 · CONEVAL 2020 · SIGEH Hidalgo · PDU Tizayuca")
badge_version(slide)

# ═══════════════════════════════════════════════════════════
# EXPORTAR
# ═══════════════════════════════════════════════════════════
OUTPUT = "/content/Tizayuca_Unidad1_v1.pptx"
prs.save(OUTPUT)
print(f"\n✅  Presentación generada: {OUTPUT}")
print(f"   7 slides · Sistema visual Industrial Organic")

from google.colab import files
files.download(OUTPUT)
